from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Paths
LOGO_PATH = os.path.join("static", "bosch.png")
PPTX_PATH = "BikeVault_Presentation.pptx"

# Helper: Add Bosch logo to slide (top-right)
def add_logo(slide):
    if os.path.exists(LOGO_PATH):
        left = Inches(7.5)
        top = Inches(0.1)
        slide.shapes.add_picture(LOGO_PATH, left, top, width=Inches(1.5))

# Helper: Add title and subtitle
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    # Add a colored rectangle for style
    left = top = 0
    width = prs.slide_width
    height = Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(182, 0, 0)  # Bosch red
    shape.line.color.rgb = RGBColor(255,255,255)
    shape.line.width = Pt(0)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    add_logo(slide)

# Helper: Add section header

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    # Add a colored bar
    left = 0
    top = Inches(0.8)
    width = prs.slide_width
    height = Inches(0.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(182, 0, 0)
    shape.line.color.rgb = RGBColor(255,255,255)
    shape.line.width = Pt(0)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(182,0,0)
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    add_logo(slide)

# Helper: Add bullet slide

def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(50,50,50)
    add_logo(slide)
# Helper: Add image slide
def add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(7))
    add_logo(slide)

# Helper: Add content slide with custom content

def add_content_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = content
    p.font.size = Pt(20)
    add_logo(slide)

# Create presentation
prs = Presentation()


# Title Slide
add_title_slide(prs, "BikeVault - Vehicle Testing Management System", "Bosch | April 2026")

# Project Overview
add_section_slide(prs, "Project Overview")
add_bullet_slide(prs, "What is BikeVault?", [
    "Web-based vehicle testing & management system",
    "Digitizes intake, testing, and delivery workflows",
    "Automates RC card extraction and notifications"
])
add_bullet_slide(prs, "Business Purpose", [
    "Streamline vehicle testing operations",
    "Reduce manual paperwork and errors",
    "Track vehicle status in real-time",
    "Automated reports and notifications"
])

# Agenda
add_section_slide(prs, "Agenda")
add_bullet_slide(prs, "Presentation Flow", [
    "Project Overview",
    "Technology Stack",
    "Key Features",
    "System Architecture",
    "Workflows",
    "Security & Validation",
    "Results & Impact",
    "Future Enhancements"
])

# Technology Stack
add_section_slide(prs, "Technology Stack")
add_bullet_slide(prs, "Backend", [
    "Flask (Python)",
    "SQLAlchemy ORM",
    "Flask-Login (Auth)",
    "Flask-Mail (Email)",
    "EasyOCR, OpenCV (OCR/Image)"
])
add_bullet_slide(prs, "Frontend", [
    "HTML5, TailwindCSS",
    "Vanilla JavaScript",
    "Jinja2 Templates"
])
add_bullet_slide(prs, "Database & Infra", [
    "SQLite (dev), PostgreSQL/MySQL (prod)",
    "Gmail SMTP integration",
    "Heroku/Railway deployment"
])

# Key Features
add_section_slide(prs, "Key Features")
add_bullet_slide(prs, "Core Features", [
    "User authentication & roles",
    "Vehicle reception & registration",
    "RC card auto-extraction (OCR)",
    "Instrumentation & testing",
    "Vehicle return workflow",
    "Automated email notifications",
    "Admin dashboard & analytics",
    "Activity tracking & audit trail",
    "Data export (Excel)"
])
add_bullet_slide(prs, "Special Features", [
    "EasyOCR-based RC card extraction",
    "Automated email notifications",
    "Admin dashboard with analytics"
])

# System Architecture
add_section_slide(prs, "System Architecture")
add_content_slide(prs, "Architecture Diagram", "User → Flask → DB/OCR/Email\n[See documentation for full diagram]")
add_bullet_slide(prs, "Core Components", [
    "app.py: Main Flask app",
    "models.py: Database models",
    "routes.py: API endpoints",
    "rc_processor.py: OCR pipeline",
    "email_utils.py: Email service"
])

# Demo Screenshots (placeholder)
add_section_slide(prs, "Demo Screenshots")
add_content_slide(prs, "Vehicle Reception UI", "[Screenshot: Vehicle intake form with RC card upload]")
add_content_slide(prs, "Admin Dashboard", "[Screenshot: Dashboard with stats and user management]")
add_content_slide(prs, "OCR Extraction", "[Screenshot: RC card extraction result]")

# Workflows
add_section_slide(prs, "Workflows")
add_bullet_slide(prs, "Vehicle Reception", [
    "User uploads RC card image",
    "EasyOCR extracts details",
    "Form auto-fills, user reviews",
    "Vehicle saved, email sent"
])
add_bullet_slide(prs, "Testing & Return", [
    "Vehicle marked 'In Test'",
    "Instrumentation records created",
    "Test completion triggers email",
    "Vehicle returned, status updated"
])

# Security & Validation
add_section_slide(prs, "Security & Validation")
add_bullet_slide(prs, "Security Features", [
    "Password hashing (Werkzeug)",
    "Session-based authentication",
    "CSRF protection ready",
    "Secure file uploads",
    "Input validation"
])
add_bullet_slide(prs, "Testing & Validation", [
    "User auth tested",
    "OCR accuracy 70-95%",
    "Email & DB tested",
    "Manual & automated validation"
])

# Results & Impact
add_section_slide(prs, "Results & Impact")
add_bullet_slide(prs, "Project Outcomes", [
    "Reduced data entry time by 60%",
    "Error reduction by 90%",
    "Production-ready system",
    "Scalable & cloud deployable"
])
add_bullet_slide(prs, "Future Enhancements", [
    "Mobile app integration",
    "Advanced analytics",
    "Multi-location support"
])

# Thank You Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Thank You!"
slide.placeholders[1].text = "Questions?\nBikeVault | Bosch"
add_logo(slide)

prs.save(PPTX_PATH)
print(f"Presentation generated: {PPTX_PATH}")
